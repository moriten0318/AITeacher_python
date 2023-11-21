import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import json

# プレゼンテーションを作成
prs = Presentation()

# 背景に使用する画像のパス
background_img_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'image//kokuban.png')

slidetext="＜関節と骨の位置の視覚化＞・関節は赤シール・骨は青色線"
# 句読点「。」で分割
sentences = slidetext.split('。')

# 文を句読点「。」で区切り、文末に「\n」を追加し、18文字ごとにさらに「\n」を追加
formatted_text = ''
for sentence in sentences:
    if sentence:
        #sentence += '。'  # 句読点を追加
        if len(sentence) > 18:
            addtext=""
            char_list = [char for char in sentence]
            for i in char_list:
                addtext += i
                if len(addtext)==18:
                    addtext+="\n"
            formatted_text +=addtext
        else:
            formatted_text += sentence
    print(formatted_text)


for i in range(5):  # 5つのスライドを作成する例
    # 新しいスライドを追加
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # 空のレイアウト

    # 背景画像を追加
    # 背景画像を拡大して追加
    width, height = prs.slide_width, prs.slide_height
    enlarged_width, enlarged_height = width * 1.03, height * 1.03
    left = (width - enlarged_width)/2  # 中央に配置
    top = (height - enlarged_height)/2

    pic = slide.shapes.add_picture(background_img_path, left, top, enlarged_width, enlarged_height)

    # スライドにテキストを追加
    textbox_left = Inches(0.5)  # テキストボックスの左の位置
    textbox_top = Inches(0.5)  # テキストボックスの上の位置（スライドの中央付近）
    textbox_width = width - Inches(2)  # テキストボックスの幅
    textbox_height = Inches(1)  # テキストボックスの高さ
    txBox = slide.shapes.add_textbox(textbox_left, textbox_top, textbox_width, textbox_height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.word_wrap = True
    p.text = formatted_text.replace("。", "。\n")
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(255, 255, 255)  # 白色

# プレゼンテーションを保存
prs.save('generated_presentation.pptx')
print("終了")
