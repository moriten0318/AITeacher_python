import os
import openai
from dotenv import load_dotenv
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

import glob
from comtypes import client
import datetime

jsonfile_name="honji_tenkai_1202.json"#渡すJSONファイル名
OUT_dir="lessondata"
# 今日の日付を取得
today = datetime.date.today()




def openjson(filename):
    with open("json\\" + filename, encoding='utf-8') as file:
    
        return json.load(file)
    
def savejson(data, OUT):
    with open(OUT, 'w', encoding='utf-8') as file:
        json.dump(data, file,ensure_ascii=False, indent=4)  # 'indent'オプションは、読みやすい形式で保存するためのものです

def text_format(slidetext):
    # 句読点「。」で分割
    sentences = slidetext.split('。')


    # 文を句読点「。」で区切り、文末に「\n」を追加し、15文字ごとにさらに「\n」を追加
    formatted_text = ''
    for sentence in sentences:
        if sentence:
            #sentence += '。'  # 句読点を追加
            char=len(sentence)
            sentence=sentence.replace("＞", "＞\n\n")
            sentence=sentence.replace("・", "\n・")

            if "\n" not in sentence and char > 18:
                addtext=""
                char_list = [char for char in sentence]
                for i in char_list:
                    addtext += i
                    if len(addtext)==18:
                        addtext+="\n"
                formatted_text +=addtext
            else:
                formatted_text += sentence
        char = 0
    return formatted_text

def create_presentation(lesson_data):
    prs = Presentation()

    # 背景に使用する画像のパス
    background_img_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'image//kokuban.png')
    # 新しい空のスライド（slide0）を追加
    slide0 = prs.slides.add_slide(prs.slide_layouts[5])  # 空のレイアウト
    # 背景画像を追加
    width, height = prs.slide_width, prs.slide_height
    enlarged_width, enlarged_height = width * 1.03, height * 1.03
    left = (width - enlarged_width)/2  # 中央に配置
    top = (height - enlarged_height)/2
    slide0.shapes.add_picture(background_img_path, left, top, enlarged_width, enlarged_height)
    
    for section in lesson_data:
        if '板書' in section:
            # 新しいスライドを追加
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # 空のレイアウト

            # 背景画像を追加
            # 背景画像を拡大して追加
            width, height = prs.slide_width, prs.slide_height
            enlarged_width, enlarged_height = width * 1.03, height * 1.03
            left = (width - enlarged_width)/2  # 中央に配置
            top = (height - enlarged_height)/2

            slide.shapes.add_picture(background_img_path, left, top, enlarged_width, enlarged_height)

            # スライドにテキストを追加
            textbox_left = Inches(0.5)  # テキストボックスの左の位置
            textbox_top = Inches(0.5)  # テキストボックスの上の位置（スライドの中央付近）
            textbox_width = width - Inches(2)  # テキストボックスの幅
            textbox_height = Inches(1)  # テキストボックスの高さ
            txBox = slide.shapes.add_textbox(textbox_left, textbox_top, textbox_width, textbox_height)
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.word_wrap = True
            formatted_text=text_format(section['板書'])
            p.text = formatted_text.replace("。", "。\n")
            p.font.size = Pt(28)
            p.font.color.rgb = RGBColor(255, 255, 255)  # 白色

            p.font.size = Pt(28)

    pptx_name=jsonfile_name.replace(".json", "")
    prs.save(f'{OUT_dir}\\{pptx_name}_{today}.pptx')
    print("プレゼンテーション作成終了")
    return f'{OUT_dir}\\{pptx_name}_{today}.pptx'


def export_img(fname, odir):
    application = client.CreateObject("Powerpoint.Application")
    application.Visible = True
    current_folder = os.getcwd()

    presentation = application.Presentations.open(os.path.join(current_folder, fname))

    export_path = os.path.join(current_folder, odir)
    presentation.Export(export_path, FilterName="png")

    presentation.close()
    application.quit()

def rename_img(odir):
    file_list = glob.glob(os.path.join(odir, "*.PNG"))
    for fname in file_list:
        new_fname = fname.replace('スライド', 'slide').lower()
        os.rename(fname, new_fname)
    print("画像出力完了")

def main():
    lesson_data = openjson(jsonfile_name)
    savejson(lesson_data,OUT_dir+"\\"+jsonfile_name)
    ppt_name=create_presentation(lesson_data)
    export_img(ppt_name,OUT_dir)
    rename_img(OUT_dir)

if __name__ == "__main__":
    main()
