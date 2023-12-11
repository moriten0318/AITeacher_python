import os
import openai
from dotenv import load_dotenv
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

import glob
from comtypes import client
import datetime

#このファイル1つで完結する
#ChatGPT4が生成したJSONファイルを読み込んで、【教師の発話]をJSONに追加できる
#ChatGPT4が生成したJSONファイルを読み込んで、[パワポのスライドを作成]できる

#Web版chatGPTに出力させたJSONファイル名を入力
json_name="honji_tenkai_1121.json"
OUT_dir="lessondata"
# 今日の日付を取得
today = datetime.date.today()


#APIキー取得して渡す
load_dotenv()
api_key = os.environ['OPENAI_API_KEY']
print(api_key)
openai.api_key=api_key


operating =\
"#あなたは小学校の教師です。\
#あなたに与えた「時間」、「学習活動」、「指導上の留意点」、「評価の観点」の項目を踏まえて、「そのセクションにおける教師の発話」を答えてください。\
#「学習活動」とは、授業内で生徒が達成するべき活動であり、教師の発話は学習活動が全て行えるように促すものになります。\
\
#出力する「教師の発話」とは、そのセクション内で行う学習活動や指導上の留意点の項目がすべて達成できるように、教師が生徒の前で話すような発言のことを指します。\
「教師の発話」の内容は、学習活動の内容に対応したものにしてください。例えば、学習活動で生徒が新たな知識を学ぶ場合は単語の説明をし、実験活動を行う場合は、実験の指示を出す必要があります。\
「教師の発話」では、教師は授業を聞いている生徒のことを意識し、常に生徒たちの学習状況を気にかける必要があります。\
例えば、「どうですか？理解できましたか？」「一度整理のために時間を取りますね。」などのように、時々教師が生徒に寄り添うような発言を入れてください。\
また、あなたは生徒に好かれる教師なので、ユーモアに富んだ語りをします。\
#教師の発話では口調を統一してください。以下の例で挙げている文章と同じ口調で話してください。\
#教師の発話以外の内容は発言しないでください。\
\
#以下は与えられる学習活動、指導上の留意点、評価の観点及び、それから生成する「教師の発話」の例です。\
時間:展開,\
学習活動:3注射器に閉じ込めた空気と水をそれぞれおして、体積と手応えの変化を調べ発表する。,\
指導上の留意点:うまく実験が行えない児童がいないか見て回る。対話が自然に生まれるように促す。,\
評価の観点:閉じ込めた空気や水を比較しながら力を加えると、手ごたえと体積に違いがあることをとらえることができたか。\
「みなさん、これから空気と水を注射器に閉じ込めてピストンを押すと、手ごたえはどうなるのかについて調べてみましょう。\
みなさんの予想はどのようなものですか？実際の結果はどのような物になるでしょうか？\
それでは実際に試してみましょう。時間を取るので、周りの人と協力して実験してみましょう。\
もしわからないことがあれば、遠慮なく質問してくださいね。」"



def generate_text(prompt, conversation_history):

    # プロンプトを会話履歴に追加
    conversation_history.append({"role": "user", "content": prompt})

    # GPT-4モデルを使用する場合
    response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo",#gpt-3.5-turbo/gpt-4
        messages=conversation_history
    )
    message = ""

    for choice in response.choices:
        message += choice.message['content']

    # 応答文を会話履歴に追加
    conversation_history.append({"role": "assistant", "content": message})
    return message

def openjson(filename):
    with open("json\\"+filename,encoding='utf-8') as file:
        return (json.load(file))
    
def speech_generate(jsonfilename,conversation_history):

    lesson_data=openjson(jsonfilename)

    for section in lesson_data:
        print(section['時間'])
        print(f"学習活動: {section['学習活動']}")
        print(f"指導上の留意点: {section['指導上の留意点']}")
        評価の観点 = section['評価の観点'] if section['評価の観点'] is not None else "なし"
        print(f"評価の観点: {評価の観点}")

        teacher_speech = generate_text("学習活動:"+section["学習活動"]+"指導上の留意点:"+section["指導上の留意点"],conversation_history)
        section["教師の発話"] = teacher_speech
        print("教師の発話"+section["教師の発話"])
        print(f"板書: {section['板書']}")

        print("\n")  # セクション間に空行を挿入

    # 新しいJSONデータとして保存
    with open(f'json\\updated_{jsonfilename}', 'w', encoding='utf-8') as file:
        json.dump(lesson_data, file, ensure_ascii=False, indent=2)   
    print("処理/保存終了です")
    return lesson_data

def text_format(slidetext):
    # 句読点「。」で分割
    sentences = slidetext.split('。')

    # 文を句読点「。」で区切り、文末に「\n」を追加し、15文字ごとにさらに「\n」を追加
    formatted_text = ''
    for sentence in sentences:
        if sentence:
            #sentence += '。'  # 句読点を追加
            sentence=sentence.replace("＞", "＞\n\n")
            sentence=sentence.replace("・", "\n・")

            if len(sentence) > 18:
                addtext=""
                char_list = [char for char in sentence]
                for i in char_list:
                    addtext += i
                    if len(addtext)==18:
                        addtext+="\n"
                formatted_text +=addtext
            else:
                formatted_text += sentence
    return formatted_text

def create_presentation(lesson_data):
    prs = Presentation()

    # 背景に使用する画像のパス
    background_img_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'image//kokuban.png')

    for section in lesson_data:
        if '板書' in section:
            # 新しいスライドを追加
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # 空のレイアウト

            # 背景画像を追加
            # 背景画像を拡大して追加
            width, height = prs.slide_width, prs.slide_height
            enlarged_width, enlarged_height = width * 1.03, height * 1.03
            left = (width - enlarged_width)/2  # 中央に配置
            top = (height - enlarged_height)/2

            slide.shapes.add_picture(background_img_path, left, top, enlarged_width, enlarged_height)

            # スライドにテキストを追加
            textbox_left = Inches(0.5)  # テキストボックスの左の位置
            textbox_top = Inches(0.5)  # テキストボックスの上の位置（スライドの中央付近）
            textbox_width = width - Inches(2)  # テキストボックスの幅
            textbox_height = Inches(1)  # テキストボックスの高さ
            txBox = slide.shapes.add_textbox(textbox_left, textbox_top, textbox_width, textbox_height)
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.word_wrap = True
            formatted_text=text_format(section['板書'])
            p.text = formatted_text.replace("。", "。\n")
            p.font.size = Pt(28)
            p.font.color.rgb = RGBColor(255, 255, 255)  # 白色

            p.font.size = Pt(28)

    pptx_name=json_name.replace(".json", "")
    prs.save(f'powerpoint\\{pptx_name}.pptx')
    print("プレゼンテーション作成終了")
    return f'powerpoint\\{pptx_name}.pptx'


def export_img(fname, odir):
    application = client.CreateObject("Powerpoint.Application")
    application.Visible = True
    current_folder = os.getcwd()

    presentation = application.Presentations.open(os.path.join(current_folder, fname))

    export_path = os.path.join(current_folder, odir)
    presentation.Export(export_path, FilterName="png")

    presentation.close()
    application.quit()

def rename_img(odir):
    file_list = glob.glob(os.path.join(odir, "*.PNG"))
    for fname in file_list:
        new_fname = fname.replace('スライド', 'slide').lower()
        os.rename(fname, new_fname)
    print("画像出力完了")


def main():
    # 会話履歴を格納するためのリストを初期化
    conversation_history = []
    #指示を追加
    conversation_history.append({"role": "system", "content": operating})
    lesson_data =speech_generate(json_name,conversation_history)#教師の発話を生成する→「updated+JSONファイルの名前で保存」
    ppt_name=create_presentation(lesson_data)#パワポの作成関数→JSONファイルの名前.
    export_img(ppt_name,"powerpoint")
    rename_img("powerpoint")

if __name__ == "__main__":

    main()